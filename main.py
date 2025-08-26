import os
import io
import pptx
from pptx.util import Pt
from flask import Flask, render_template, jsonify, request, send_file
from langchain_core.prompts import ChatPromptTemplate
from langchain_google_genai import ChatGoogleGenerativeAI
from dotenv import load_dotenv
from pptx.enum.text import MSO_AUTO_SIZE

# --- 1. Setup and Configuration ---
load_dotenv()
FALLBACK_GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
MODEL_NAME = "gemini-1.5-pro-latest"

# --- 2. AI Chains for Different Tasks ---

# Prompt for the initial plan generation
generate_prompt = ChatPromptTemplate.from_messages([
    ("system", "You are an expert presentation designer AI. Your goal is to generate a detailed markdown string for a presentation based on the user's text. You must structure the text into the exact number of slides requested. Your final output must ONLY be the markdown string."),
    ("user", """
    Here are the precise formatting rules you MUST follow for each slide.

    **Slide 1: Title Slide**
    - A title line starting with `#`.
    - A subtitle line starting with `##`.
    ---
    **Content Slides**
    - Each slide MUST have a title line starting with `#`.
    - Each slide MUST have 2-5 bullet points starting with `-`.
    ---
    **Conclusion Slide**
    - The final slide MUST have a title line starting with `#`.
    - The final slide MUST have 2-3 summary bullet points.

    **Optional Guidance for tone and structure:** {guidance}

    Now, analyze the following text and generate the markdown for a presentation with exactly {number_of_slides} slides.

    <TEXT>
    {bulk_text}
    </TEXT>
    """),
])

# Prompt for the "Improvise" feature
improvise_prompt = ChatPromptTemplate.from_messages([
    ("system", "You are an expert presentation editor AI. Your task is to revise and improve an existing presentation plan based on user feedback. Refine the content, titles, and structure to better match the user's guidance. Your final output must ONLY be the revised markdown string, keeping the same number of slides and the same `---` format."),
    ("user", """
    Please improvise the following presentation plan.

    **User's new guidance:** "{guidance}"

    **Original Markdown Plan:**
    ---
    {markdown_plan}
    ---
    """)
])


# --- 3. PowerPoint Helper Function ---
def create_ppt_with_template(markdown_slides: str, template_file=None):
    """Creates a PowerPoint presentation, using a template file if provided, or a new one if not."""
    if template_file:
        prs = pptx.Presentation(template_file)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        prs = pptx.Presentation()

    slides_content = [s.strip() for s in markdown_slides.strip().split('---') if s.strip()]
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    for i, slide_markdown in enumerate(slides_content):
        lines = [line.strip() for line in slide_markdown.split('\n') if line.strip()]
        if not lines: continue

        if i == 0:
            slide = prs.slides.add_slide(title_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            title.text = lines[0].replace('#', '').strip()
            if subtitle and len(lines) > 1:
                subtitle.text = lines[1].replace('##', '').strip()
        else:
            slide = prs.slides.add_slide(content_layout)
            title = slide.shapes.title
            body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
            title.text = lines[0].replace('#', '').strip()
            if body:
                tf = body.text_frame
                tf.clear()
                # --- CRITICAL FIX: Properly populate bullet points ---
                # The first paragraph is set outside the loop to establish the text frame.
                # Subsequent paragraphs are added correctly.
                first_bullet = True
                for line in lines[1:]:
                    if line.startswith('-'):
                        if first_bullet:
                            p = tf.paragraphs[0]
                            first_bullet = False
                        else:
                            p = tf.add_paragraph()
                        p.text = line.lstrip('- ').strip()
                        p.font.size = Pt(18)
                        p.level = 0
                # Enable auto-fit after all text is added
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# --- 4. Flask Web Application ---
app = Flask(__name__)

def get_llm(user_api_key):
    api_key_to_use = user_api_key or FALLBACK_GOOGLE_API_KEY
    if not api_key_to_use:
        raise ValueError("API key is missing.")
    return ChatGoogleGenerativeAI(model=MODEL_NAME, google_api_key=api_key_to_use)

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/generate_plan", methods=['POST'])
def generate_plan_route():
    data = request.get_json()
    try:
        llm = get_llm(data.get("api_key"))
        chain = generate_prompt | llm
        response = chain.invoke({
            "bulk_text": data["bulk_text"],
            "guidance": data.get("guidance", "A standard professional presentation."),
            "number_of_slides": data["number_of_slides"]
        })
        return jsonify({"markdown_plan": response.content.strip().replace("```markdown", "").replace("```", "")})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/improvise_plan", methods=['POST'])
def improvise_plan_route():
    data = request.get_json()
    try:
        llm = get_llm(data.get("api_key"))
        chain = improvise_prompt | llm
        response = chain.invoke({
            "markdown_plan": data["markdown_plan"],
            "guidance": data.get("guidance", "Make it more professional.")
        })
        return jsonify({"markdown_plan": response.content.strip().replace("```markdown", "").replace("```", "")})
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/create_file", methods=['POST'])
def create_file_route():
    template_file = None
    if 'template_file' in request.files and request.files['template_file'].filename != '':
        template_file = request.files['template_file']
    
    markdown_plan = request.form.get('markdown_plan')
    if not markdown_plan:
        return jsonify({"error": "Missing markdown plan."}), 400
    
    try:
        ppt_buffer = create_ppt_with_template(markdown_plan, template_file)
        return send_file(
            ppt_buffer,
            as_attachment=True,
            download_name="Generated-Presentation.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)
